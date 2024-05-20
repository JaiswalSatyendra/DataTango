from pptx import Presentation
from pptx.util import Inches
import os
import pandas as pd
from langchain.agents import AgentType
from langchain_experimental.agents import create_pandas_dataframe_agent
from langchain.chat_models import ChatOpenAI
from dotenv import load_dotenv
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
import io

# Load environment variables from .env file
load_dotenv()

# Supported file formats for loading data
file_formats = {
    "csv": pd.read_csv,
    "xls": pd.read_excel,
    "xlsx": pd.read_excel,
    "xlsm": pd.read_excel,
    "xlsb": pd.read_excel,
}

# Function to load data from provided file path
def load_data(file_path):
    if file_path.endswith(".csv"):
        return pd.read_csv(file_path)
    else:
        print(f"Unsupported file format: {os.path.splitext(file_path)[1]}")
        return None

# Dataset path (Specify the path to your dataset file)
dataset_path = "ccg_survey_continued_amber.csv"

# Create a presentation object
prs = Presentation()

# Load data
df = load_data(dataset_path)

# Retrieve OpenAI API key from environment variables
openai_api_key = os.getenv("OPENAI_API_KEY")

# Check if OpenAI API key is available
if not openai_api_key:
    print("OpenAI API key is missing. Please make sure it's set in your environment variables.")
    exit()

# Initialize ChatOpenAI instance
llm = ChatOpenAI(
    temperature=0, model="gpt-4-turbo", openai_api_key=openai_api_key, streaming=True
)

# Create agent for pandas DataFrame
pandas_df_agent = create_pandas_dataframe_agent(
    llm,
    df,
    verbose=True,
    agent_type=AgentType.OPENAI_FUNCTIONS,
    handle_parsing_errors=True,
)

# Ask a question to OpenAI
prompt = input("Enter your question: ")
messages = [{"role": "user", "content": prompt}]
response = pandas_df_agent.run(messages)

# Generate a plot based on the response
# Assuming response is a list of values for plotting
plt.figure()
plt.bar(range(len(response)), response)
plt.xlabel('X-axis')
plt.ylabel('Y-axis')
plt.title('Bar Chart Title')

# Save the plot as an image
image_stream = io.BytesIO()
plt.savefig(image_stream, format='png')
image_stream.seek(0)

# Add a slide with a title and content
slide_layout = prs.slide_layouts[5]  # Use layout index 5 for a blank slide
slide = prs.slides.add_slide(slide_layout)

# Add the plot image to the slide
left_inch = Inches(1)
top_inch = Inches(2)
slide.shapes.add_picture(image_stream, left_inch, top_inch, width=Inches(6), height=Inches(4.5))

# Save the presentation
prs.save("sample.pptx")
